#!/usr/bin/env python3
"""KHOJ — fill the INNOHACK template into the finished pitch deck.

Run AFTER the template has been copied to KHOJ_INNOHACK.pptx and expanded to
16 slides with add_slide.py.

    py build_deck.py

It keeps every template design element (logos, watermark, footer, the navy
INNOHACK/section headings) and replaces only the body placeholders, then lays
out its own text columns, status cards and diagrams on top.
"""
import os
import shutil
import subprocess
import sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, "KHOJ_INNOHACK.pptx")
TEMPLATE = r"C:\Users\rishi\Downloads\INNOHACK 2.0 TEMPLATE.pptx"
ADD_SLIDE = (r"C:\Users\rishi\AppData\Roaming\Claude\local-agent-mode-sessions"
             r"\skills-plugin\c1e1c438-10de-421c-86df-d9b8f3b2aec2"
             r"\5c287bc0-fab1-43aa-9272-00c1f3258702\skills\pptx\scripts\add_slide.py")
EXTRA_SLIDES = 6          # 10 template slides + 6 duplicates = 16


def prepare():
    """Rebuild from the pristine template every run.

    This MUST happen each time: the builder adds textboxes and pictures, so
    running it twice against an already-built deck would stack a second copy of
    every shape on top of the first. Starting from the template makes the build
    idempotent — you can iterate on content freely.
    """
    if not os.path.exists(TEMPLATE):
        sys.exit("Template not found: " + TEMPLATE)
    shutil.copyfile(TEMPLATE, DECK)
    for _ in range(EXTRA_SLIDES):
        r = subprocess.run([sys.executable, ADD_SLIDE, DECK, "slide4.xml",
                            "--after", "slide4.xml"],
                           capture_output=True, text=True)
        if r.returncode != 0:
            sys.exit("add_slide failed:\n" + r.stdout + r.stderr)

# ---- canvas -----------------------------------------------------------------
W, H = 18288000, 10287000
L = 886241                      # template's left margin
R = W - L
Y0, Y1 = 3_350_000, 9_640_000   # usable band between heading and footer
COLW = 6_900_000                # left text column
IMGX, IMGW = 8_000_000, 9_400_000

# ---- palette (matches the diagrams) -----------------------------------------
NAVY  = RGBColor(0x2E, 0x31, 0x92)   # template heading navy
INK   = RGBColor(0x1B, 0x2A, 0x3A)
MUTED = RGBColor(0x5A, 0x6B, 0x7C)
CYAN  = RGBColor(0x0E, 0x8F, 0xA0)
AMBER = RGBColor(0xD9, 0x88, 0x29)
RED   = RGBColor(0xB8, 0x43, 0x3A)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT  = {CYAN: RGBColor(0xE4, 0xF3, 0xF5), AMBER: RGBColor(0xFC, 0xF1, 0xE2),
         RED: RGBColor(0xF9, 0xE7, 0xE5), GREEN: RGBColor(0xE4, 0xF1, 0xEA),
         NAVY: RGBColor(0xEC, 0xEE, 0xF7)}

FONT = "Calibri"
PLACEHOLDER_TOPS = {3545509, 4619653, 5700684, 6836797}
TITLE_TOP = 2216419


# ---- template surgery -------------------------------------------------------
def body_shapes(slide):
    return [s for s in slide.shapes
            if any(abs(s.top - t) < 20000 for t in PLACEHOLDER_TOPS)]


def strip_body(slide):
    """Remove the template's body placeholders, leaving all design intact."""
    for s in body_shapes(slide):
        s._element.getparent().remove(s._element)


def set_heading(slide, text):
    """Rewrite the navy section heading, preserving its run formatting."""
    for s in slide.shapes:
        if abs(s.top - TITLE_TOP) < 20000 and s.has_text_frame:
            p = s.text_frame.paragraphs[0]
            for r in list(p.runs)[1:]:
                r._r.getparent().remove(r._r)
            if p.runs:
                p.runs[0].text = text
            return


def set_placeholder(slide, top, text, size=None, color=None, bold=None):
    for s in body_shapes(slide):
        if abs(s.top - top) < 20000 and s.has_text_frame:
            p = s.text_frame.paragraphs[0]
            for r in list(p.runs)[1:]:
                r._r.getparent().remove(r._r)
            if not p.runs:
                return
            r = p.runs[0]
            r.text = text
            if size:  r.font.size = Pt(size)
            if color: r.font.color.rgb = color
            if bold is not None: r.font.bold = bold
            return


# ---- content builders -------------------------------------------------------
def textbox(slide, x, y, w, h, blocks):
    """blocks: list of dicts {t, size, bold, color, space, bullet, italic}"""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(b.get("space", 8))
        if b.get("bullet"):
            rb = p.add_run()
            rb.text = "▪   "
            rb.font.size = Pt(b.get("size", 19))
            rb.font.color.rgb = b.get("bcolor", CYAN)
            rb.font.bold = True
            rb.font.name = FONT
        r = p.add_run()
        r.text = b["t"]
        r.font.size = Pt(b.get("size", 19))
        r.font.bold = b.get("bold", False)
        r.font.italic = b.get("italic", False)
        r.font.color.rgb = b.get("color", INK)
        r.font.name = FONT
    return tb


def card(slide, x, y, w, h, big, label, sub, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = TINT[color]
    sh.line.color.rgb = color
    sh.line.width = Pt(1.5)
    sh.adjustments[0] = 0.08
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(220000)
    tf.margin_right = Emu(140000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT        # shapes default to centred; text reads left
    p.space_after = Pt(2)
    r = p.add_run(); r.text = big
    r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
    r2 = p.add_run(); r2.text = "   " + label
    r2.font.size = Pt(19); r2.font.bold = True; r2.font.color.rgb = INK; r2.font.name = FONT
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r3 = p2.add_run(); r3.text = sub
        r3.font.size = Pt(15); r3.font.color.rgb = MUTED; r3.font.name = FONT
    return sh


def cards(slide, items, x=IMGX, w=IMGW, y0=Y0, gap=190000):
    n = len(items)
    h = int(((Y1 - y0) - gap * (n - 1)) / n)
    for i, (big, lab, sub, col) in enumerate(items):
        card(slide, x, y0 + i * (h + gap), w, h, big, lab, sub, col)


def picture(slide, name, x, y, w):
    p = os.path.join(HERE, name)
    iw, ih = Image.open(p).size
    h = int(w * ih / iw)
    slide.shapes.add_picture(p, Emu(x), Emu(y), Emu(w), Emu(h))
    return h


def pic_fit(slide, name, x, y, maxw, maxh):
    """Fit inside a box, centred horizontally in that box."""
    p = os.path.join(HERE, name)
    iw, ih = Image.open(p).size
    w = maxw
    h = int(w * ih / iw)
    if h > maxh:
        h = maxh
        w = int(h * iw / ih)
    slide.shapes.add_picture(p, Emu(int(x + (maxw - w) / 2)), Emu(int(y + (maxh - h) / 2)),
                             Emu(w), Emu(h))


def diagram_slide(slide, heading, takeaway, image):
    """Heading + a one-line takeaway + the diagram filling the rest.

    The takeaway matters: the slide band is far wider than a 16:9 diagram, so a
    centred image alone leaves the slide looking empty and, worse, makes a judge
    work out the point for themselves. One sentence fixes both.
    """
    strip_body(slide)
    set_heading(slide, heading)
    textbox(slide, L, 3_020_000, 16_500_000, 460_000, [
        {"t": takeaway, "size": 20, "italic": True, "color": MUTED, "space": 0},
    ])
    pic_fit(slide, image, L, 3_560_000, 16_500_000, 6_060_000)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# =============================================================================
def main():
    prepare()
    prs = Presentation(DECK)
    s = list(prs.slides)
    if len(s) != 16:
        sys.exit("expected 16 slides, got %d" % len(s))

    # ---------------------------------------------------------------- 1 TITLE
    strip_body(s[0])
    set_heading(s[0], "KHOJ")
    # NOTE: width must stop short of IMGX (8.0M) or this column runs under the
    # proof cards on the right.
    textbox(s[0], L, 3_150_000, COLW, 2_600_000, [
        {"t": "A leaderless rescue swarm\nthat finds people it cannot see.",
         "size": 34, "bold": True, "color": NAVY, "space": 14},
        {"t": "“Other swarms search together.\nKHOJ finds the invisible.”",
         "size": 21, "italic": True, "color": MUTED, "space": 0},
    ])
    textbox(s[0], L, 6_450_000, COLW, 2_900_000, [
        {"t": "TRACK", "size": 14, "bold": True, "color": CYAN, "space": 2},
        {"t": "[YOUR TRACK NAME]", "size": 22, "bold": True, "color": INK, "space": 14},
        {"t": "PROBLEM STATEMENT", "size": 14, "bold": True, "color": CYAN, "space": 2},
        {"t": "DIOT-01 — Autonomous Drone Swarm\nfor Search and Rescue Operations",
         "size": 19, "bold": True, "color": INK, "space": 4},
        {"t": "Aligned SDGs 9 · 11 · 16", "size": 15, "color": MUTED, "space": 14},
        {"t": "TEAM  [TEAM NAME]", "size": 20, "bold": True, "color": INK, "space": 4},
        {"t": "[MEMBER 1 — REG NO.]   ·   [MEMBER 2 — REG NO.]\n"
              "[MEMBER 3 — REG NO.]   ·   [MEMBER 4 — REG NO.]",
         "size": 15, "color": MUTED, "space": 0},
    ])
    # a proof strip on the title slide: hard numbers before the first claim
    cards(s[0], [
        ("0.0%", "mesh packet loss",
         "Measured across five ESP32 boards on real hardware.", CYAN),
        ("< 2 s", "to detect a dead drone",
         "Every surviving board concludes it independently — no leader.", RED),
        ("4-5 / 5", "survivors found in simulation",
         "A camera-only baseline manages 3-4, and never the hidden one.", GREEN),
    ], y0=3_150_000)
    notes(s[0], "One line: a leaderless drone swarm that finds victims cameras cannot see, "
                "by cooperatively locating their phone's radio signal.")

    # ------------------------------------------------------------- 2 PROBLEM
    strip_body(s[1])
    set_heading(s[1], "PROBLEM STATEMENT")
    textbox(s[1], L, Y0, COLW, 5_800_000, [
        {"t": "When a building collapses, the victims who most need finding are "
              "exactly the ones nobody can see — buried under slab, trapped in "
              "voids, unconscious and silent.", "size": 20, "space": 14},
        {"t": "Search-and-rescue drones today are camera-based. If the camera "
              "cannot see you, the drone cannot find you.",
         "size": 20, "bold": True, "color": NAVY, "space": 20},
        {"t": "Why it matters", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "The golden 72 hours — survival falls steeply with every hour "
              "under rubble. Speed is the whole outcome.", "size": 18, "bullet": True, "space": 9},
        {"t": "Disaster zones have no infrastructure by definition: no network, "
              "no guaranteed GPS, no reliable power.", "size": 18, "bullet": True, "space": 9},
        {"t": "A false positive costs a rescue team three minutes. A false negative "
              "costs a life. No current system treats those differently.",
         "size": 18, "bullet": True, "space": 9},
    ])
    cards(s[1], [
        ("01", "Blind to the buried",
         "Aerial cameras find only visible victims. Occluded and half-buried "
         "survivors are invisible to the entire fleet.", RED),
        ("02", "Fragile coordination",
         "One ground station computes every trajectory and needs RTK-GPS plus a live "
         "uplink — the first things a disaster destroys.", AMBER),
        ("03", "The threshold trap",
         "Lower the detector threshold and responders drown in false alarms. Keep it "
         "high and half-buried people are discarded as noise.", NAVY),
    ])
    notes(s[1], "Three failures. The third is the one nobody talks about: detection "
                "thresholds force a choice between missing victims and flooding responders.")

    # ------------------------------------------------------------ 3 SOLUTION
    strip_body(s[2])
    set_heading(s[2], "PROPOSED SOLUTION")
    textbox(s[2], L, Y0, COLW, 5_900_000, [
        {"t": "A leaderless drone swarm that finds people cameras cannot see — "
              "by cooperatively locating a trapped victim’s phone from its radio signal.",
         "size": 20, "space": 18},
        {"t": "1.  Cooperative RF localization — the hero",
         "size": 21, "bold": True, "color": CYAN, "space": 6},
        {"t": "Several drones share signal-strength readings over their own mesh and climb "
              "the gradient to a phone none of them can see. One reading is meaningless; "
              "four readings from four positions point at a person.",
         "size": 18, "color": INK, "space": 16},
        {"t": "2.  The re-observation loop — the spine",
         "size": 21, "bold": True, "color": AMBER, "space": 6},
        {"t": "A 0.41-confidence sighting is too strong to ignore and too weak to trust. "
              "The swarm turns that doubt into an auctioned task: another drone flies a "
              "different angle, and the two independent looks are fused into a confirm "
              "or a dismissal.", "size": 18, "color": INK, "space": 16},
    ])
    textbox(s[2], IMGX, Y0, IMGW, 5_900_000, [
        {"t": "WHAT MAKES IT INNOVATIVE", "size": 22, "bold": True, "color": NAVY, "space": 14},
        {"t": "The swarm bids in information gain, not area.",
         "size": 26, "bold": True, "color": CYAN, "space": 12},
        {"t": "Searching empty ground, re-checking an uncertain sighting, and triangulating "
              "an invisible phone all collapse into the same question:",
         "size": 18, "space": 10},
        {"t": "“Where will my next move most reduce what the swarm does not yet know?”",
         "size": 20, "italic": True, "bold": True, "color": NAVY, "space": 18},
        {"t": "That is why this is one system and not three features — and it answers "
              "the hardest question a judge can ask:", "size": 18, "space": 10},
        {"t": "“Why is this a swarm and not five drones with a shared to-do list?”",
         "size": 19, "bold": True, "color": INK, "space": 8},
        {"t": "Because a to-do list divides work. KHOJ divides doubt.",
         "size": 21, "bold": True, "color": AMBER, "space": 0},
    ])
    notes(s[2], "Lead with the hero. The judge decides whether they care in 30 seconds: "
                "a hidden phone found by drones that cannot see it.")

    # -------------------------------------------------------- 4 ARCHITECTURE
    # NOTE: each takeaway must say something the diagram's own subtitle does not.
    # Repeating it puts the same sentence twice on screen and looks careless.
    diagram_slide(s[3], "SYSTEM ARCHITECTURE",
                  "Five ESP32s, one laptop, and two links that never mix — with a real "
                  "F450 dropping in as one more agent.",
                  "01_architecture.png")
    notes(s[3], "The laptop simulates and draws but never chooses a goal. Every goal on "
                "screen was decided on a separate microcontroller and agreed by radio.")

    # ----------------------------------------------------- 5 SOFTWARE STACK
    strip_body(s[4])
    set_heading(s[4], "SOFTWARE STACK  ·  THE ALGORITHMS")
    textbox(s[4], L, Y0, COLW, 5_900_000, [
        {"t": "Runs on the laptop. Holds ground truth, simulates the world, and provides "
              "the reference implementation of every decision rule.",
         "size": 19, "space": 16},
        {"t": "Built and unit-tested", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "Belief map — 32×32 survivor-probability grid per agent, drained "
              "multiplicatively over searched ground", "size": 17, "bullet": True, "space": 8},
        {"t": "Unified auction — one bid function across all four task types",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Log-odds fusion — Bayesian confirm / dismiss on fused confidence",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Hive-Mind — broadcast dismissals so no drone repeats a ruled-out lead",
         "size": 17, "bullet": True, "space": 8},
        {"t": "World simulator — physics, RF path-loss, low-threshold detector",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Lawnmower baseline + metrics harness", "size": 17, "bullet": True, "space": 8},
        {"t": "YOLO perception on real aerial imagery, dashboard",
         "size": 17, "bullet": True, "space": 14},
        {"t": "Remaining: integration — binding the modules into one live loop "
              "and onto the serial contract.",
         "size": 18, "bold": True, "color": AMBER, "space": 0},
    ])
    cards(s[4], [
        ("Python 3", "stdlib only, zero dependencies",
         "The core runs headless on any machine — no install step to fail on demo day.", CYAN),
        ("4-5 / 5", "survivors confirmed in simulation",
         "Against an identical world, the camera-only baseline manages 3-4.", GREEN),
        ("Reference", "for the firmware port",
         "The same auction is being ported to C so the boards and the simulator "
         "provably agree.", NAVY),
    ])
    notes(s[4], "Software is feature-complete at module level; integration is the "
                "remaining work. Be precise about that distinction.")

    # ----------------------------------------------------------- 6 AUCTION
    diagram_slide(s[5], "THE AUCTION  ·  LEADERLESS CONSENSUS",
                  "Every board runs the same rule on the same data — so all five reach "
                  "the same winner with no referee.",
                  "02_auction_round.png")
    notes(s[5], "Same data plus same rule equals same winner, computed independently on "
                "every chip. There is no referee to kill.")

    # ---------------------------------------------------- 7 RE-OBSERVATION
    diagram_slide(s[6], "RE-OBSERVATION  ·  WHY IT IS A SWARM",
                  "Run the detector at a low threshold to catch half-buried victims — "
                  "then cancel the false alarms with a second look.",
                  "03_reobservation.png")
    notes(s[6], "This is the intellectual spine. It lets us run the detector at a low "
                "threshold without flooding responders with false alarms.")

    # ----------------------------------------------------- 8 HARDWARE STACK
    strip_body(s[7])
    set_heading(s[7], "HARDWARE STACK  ·  THE DECENTRALIZED MESH")
    textbox(s[7], L, Y0, COLW, 5_900_000, [
        {"t": "Runs on the boards. Five ESP32s negotiating peer-to-peer over their own "
              "radio, with no router, no pairing and no leader.",
         "size": 19, "space": 16},
        {"t": "Working on real hardware", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "ESP-NOW broadcast mesh — 23-byte packets at 5 Hz, all five boards",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Live peer table — every board independently tracks who is alive",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Failure detection — silence for 2 s and the peer is dropped everywhere",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Per-agent RSSI capture straight from the receive callback",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Byte-verified USB contract, proven both directions",
         "size": 17, "bullet": True, "space": 8},
        {"t": "MAC-derived identity — one binary flashes to every board",
         "size": 17, "bullet": True, "space": 14},
        {"t": "Remaining: port the auction on-device, then share RF samples across the mesh.",
         "size": 18, "bold": True, "color": AMBER, "space": 0},
    ])
    cards(s[7], [
        ("0.0%", "measured packet loss",
         "Across five boards on the bench — computed from gaps in each sender’s "
         "sequence numbers, not estimated.", CYAN),
        ("< 2 s", "to detect a dead drone",
         "Every surviving board reaches that conclusion independently. Recovery is "
         "emergent, not scripted.", RED),
        ("-23 to -72 dBm", "live signal range",
         "Per-agent RSSI working on hardware — the measurement the entire RF hero "
         "demo depends on.", GREEN),
    ])
    notes(s[7], "These are measured numbers from our own bench, not datasheet figures. "
                "Loss is computed from sequence-number gaps.")

    # -------------------------------------------------- 9 RF LOCALIZATION
    diagram_slide(s[8], "COOPERATIVE RF LOCALIZATION  ·  THE HERO",
                  "Gradient-seeking, not trilateration — it never needs calibration, "
                  "only “stronger here, or there?”",
                  "04_rf_localization.png")
    notes(s[8], "Gradient-seeking, not trilateration. It never needs absolute distance — "
                "only 'stronger here or there', which is what survives rubble.")

    # ------------------------------------------------ 10 WHY DECENTRALIZED
    diagram_slide(s[9], "WHY DECENTRALIZED",
                  "O(1) per drone instead of O(N) on a ground station — which is exactly "
                  "why a $3 chip is enough.",
                  "05_centralized_vs_khoj.png")
    notes(s[9], "State the trade-off before the judge finds it: this does not scale to "
                "1000 drones. For six drones with no infrastructure it is the only thing "
                "that works.")

    # ---------------------------------------------------------- 11 RESULTS
    diagram_slide(s[10], "DEMO  ·  PROTOTYPE  ·  RESULTS",
                  "Verified across four random worlds — plus the one victim a camera-only "
                  "sweep can never find.",
                  "06_results.png")
    notes(s[10], "Demo beats: self-organise, re-observation confirm, Hive-Mind dismissal, "
                 "yank a drone's power, then the hidden phone. Never cut the power-yank.")

    # ----------------------------------------------------------- 12 IMPACT
    strip_body(s[11])
    set_heading(s[11], "IMPACT & USE CASES")
    textbox(s[11], L, Y0, COLW, 5_900_000, [
        {"t": "Who benefits", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "Trapped survivors — specifically the ones every camera-based system "
              "is guaranteed to miss", "size": 18, "bullet": True, "space": 9},
        {"t": "First responders — arrive with a ranked map of where people probably "
              "are, with false alarms already filtered", "size": 18, "bullet": True, "space": 9},
        {"t": "Disaster authorities — a capability needing no network, no GPS base "
              "station, no central computer", "size": 18, "bullet": True, "space": 18},
        {"t": "Real-world scenarios", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "Earthquake and building collapse — locates phones inside voids no "
              "camera can see", "size": 18, "bullet": True, "bcolor": AMBER, "space": 9},
        {"t": "Landslide and flood rescue — works with no cellular network at all",
         "size": 18, "bullet": True, "bcolor": AMBER, "space": 9},
        {"t": "Missing person in open terrain — climbs to an emitting device across "
              "featureless area", "size": 18, "bullet": True, "bcolor": AMBER, "space": 9},
        {"t": "Comms blackout — drones reposition as relays, becoming temporary "
              "disaster infrastructure (SDG 9)", "size": 18, "bullet": True,
         "bcolor": AMBER, "space": 0},
    ])
    cards(s[11], [
        ("$3", "the decision-maker",
         "A centralized equivalent needs a ground station computing O(N) trajectories. "
         "KHOJ is O(1) per drone, in parallel.", CYAN),
        ("Zero", "infrastructure required",
         "No RTK-GPS, no uplink, no internet. It runs in the conditions disasters "
         "actually create.", GREEN),
        ("Graceful", "degradation, not failure",
         "Any drone can die at any moment and the mission continues. Adding a drone "
         "means flashing one more board.", NAVY),
    ])
    notes(s[11], "The cost argument lands hard with judges: the intelligence is a $3 chip, "
                 "not a datacentre.")

    # -------------------------------------------------------- 13 CHALLENGES
    strip_body(s[12])
    set_heading(s[12], "CHALLENGES FACED")
    textbox(s[12], L, Y0, 7_900_000, 5_900_000, [
        {"t": "1.  The radio could not measure signal strength",
         "size": 23, "bold": True, "color": RED, "space": 6},
        {"t": "The entire RF concept needs each drone reading a peer’s signal strength. "
              "The standard toolchain ships a core whose receive callback returns sender and "
              "payload — but no signal strength at all. A dedicated sniffer board could "
              "not fix it: the gradient needs each drone sampling at its own position.",
         "size": 17.5, "space": 6},
        {"t": "Solved — migrated to a newer Arduino-ESP32 core, pinned it in the build, "
              "and made the firmware print rx_rssi=YES/NO at boot so it can never regress "
              "silently. Caught in hour 1, not hour 30.",
         "size": 17.5, "bold": True, "color": GREEN, "space": 16},
        {"t": "2.  Five boards, a working radio, and total silence",
         "size": 23, "bold": True, "color": RED, "space": 6},
        {"t": "Every board transmitted and none received. Two compounding bugs: the WiFi "
              "driver returned an all-zero MAC when read too early, and our identity scheme "
              "derived board IDs from that MAC — so every board took the same ID and "
              "discarded the others’ packets as its own echo.", "size": 17.5, "space": 6},
        {"t": "Solved — read the MAC from eFuse instead of the driver, and deleted the "
              "self-filter entirely, since a broadcast is never echoed to its sender.",
         "size": 17.5, "bold": True, "color": GREEN, "space": 0},
    ])
    textbox(s[12], 9_300_000, Y0, 8_100_000, 5_900_000, [
        {"t": "3.  Silent flashing failures",
         "size": 23, "bold": True, "color": RED, "space": 6},
        {"t": "With five identical boards, flashing by hand meant believing all five were "
              "updated when only three were — discovered only through a confusing "
              "peer list.", "size": 17.5, "space": 6},
        {"t": "Solved — an automated flash script that finds every connected board and "
              "prints a pass/fail table, so a missed board is impossible to overlook.",
         "size": 17.5, "bold": True, "color": GREEN, "space": 16},
        {"t": "4.  A result we chose not to hide",
         "size": 23, "bold": True, "color": AMBER, "space": 6},
        {"t": "Our system finds more survivors than the baseline, but is not yet faster to "
              "the first one — the lawnmower sweep beelines while KHOJ self-organises. "
              "We could have quietly dropped the metric.", "size": 17.5, "space": 6},
        {"t": "Instead we report it, and we know the fix: seeding the belief map with a "
              "prior. We would rather present a real limitation than a chart a judge can "
              "take apart.", "size": 17.5, "bold": True, "color": GREEN, "space": 0},
    ])
    notes(s[12], "Challenge 1 is the strongest story: we found a silent architectural "
                 "blocker in hour one because we verified instead of assuming.")

    # ------------------------------------------------------- 14 FUTURE SCOPE
    strip_body(s[13])
    set_heading(s[13], "FUTURE SCOPE")
    textbox(s[13], L, Y0, COLW, 5_900_000, [
        {"t": "Near-term", "size": 22, "bold": True, "color": CYAN, "space": 10},
        {"t": "Port the auction fully on-device so decisions never touch a laptop",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Belief-map priors — closes the time-to-first-survivor gap",
         "size": 17, "bullet": True, "space": 8},
        {"t": "RELAY task type — drones bidding to become comms relays",
         "size": 17, "bullet": True, "space": 8},
        {"t": "Rescue-route generation so responders can reach the survivor, not just "
              "locate them", "size": 17, "bullet": True, "space": 18},
        {"t": "Medium-term", "size": 22, "bold": True, "color": AMBER, "space": 10},
        {"t": "Full F450 fleet with onboard compute and real GPS",
         "size": 17, "bullet": True, "bcolor": AMBER, "space": 8},
        {"t": "Real phone detection via dedicated sniffers instead of a planted beacon",
         "size": 17, "bullet": True, "bcolor": AMBER, "space": 8},
        {"t": "Multi-hop routing to extend coverage beyond single-radio range",
         "size": 17, "bullet": True, "bcolor": AMBER, "space": 8},
        {"t": "Hash-chained evidence log — tamper-evident record of what was found, "
              "when, and by which drone (SDG 16)", "size": 17, "bullet": True,
         "bcolor": AMBER, "space": 0},
    ])
    textbox(s[13], IMGX, Y0 + 600000, IMGW, 4_600_000, [
        {"t": "LONG-TERM VISION", "size": 22, "bold": True, "color": NAVY, "space": 16},
        {"t": "A disposable, infrastructure-free rescue layer.",
         "size": 28, "bold": True, "color": CYAN, "space": 16},
        {"t": "A crate of low-cost drones that any municipal fire service can carry to a "
              "collapse site, throw into the air, and have self-organise into a search team "
              "within seconds — no operator assigning targets, no ground station, "
              "no network.", "size": 19, "space": 16},
        {"t": "The intelligence lives in the fleet, not in a building that may itself "
              "have fallen down.", "size": 21, "bold": True, "italic": True,
         "color": NAVY, "space": 0},
    ])
    notes(s[13], "Close the future slide on the vision line — it is the most quotable "
                 "sentence in the deck.")

    # -------------------------------------------------------- 15 CONCLUSION
    strip_body(s[14])
    set_heading(s[14], "CONCLUSION")
    textbox(s[14], L, Y0, COLW, 5_900_000, [
        {"t": "What we built", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "A five-board decentralized mesh with 0.0% measured packet loss and "
              "sub-2-second failure detection", "size": 18, "bullet": True, "space": 9},
        {"t": "A unified auction where routine search, uncertain-sighting re-checks and "
              "invisible-victim localization are the same decision",
         "size": 18, "bullet": True, "space": 9},
        {"t": "A complete simulation that beats a camera-only baseline on survivors "
              "found — and finds the RF-only victim every single run, which the "
              "baseline can never do", "size": 18, "bullet": True, "space": 18},
        {"t": "Why it matters", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "Every drone swarm you have seen is centrally controlled — a ground "
              "station computing every trajectory, needing GPS and a live uplink. That "
              "architecture is excellent for a light show and useless in a collapsed "
              "building.", "size": 18, "space": 10},
        {"t": "KHOJ is built for the conditions disasters actually create: no network, "
              "no infrastructure, and no guarantee that every drone survives the hour.",
         "size": 19, "bold": True, "color": NAVY, "space": 0},
    ])
    textbox(s[14], IMGX, Y0 + 700000, IMGW, 4_400_000, [
        {"t": "THE NOVELTY", "size": 22, "bold": True, "color": NAVY, "space": 16},
        {"t": "The swarm bids in information gain, not area.",
         "size": 30, "bold": True, "color": CYAN, "space": 18},
        {"t": "Cooperative RF localization is the proof of it: several drones share what "
              "each one alone cannot interpret, and together point at a person nobody "
              "can see.", "size": 19, "space": 24},
        {"t": "“Other swarms search together.\nKHOJ finds the invisible.”",
         "size": 26, "bold": True, "italic": True, "color": NAVY, "space": 0},
    ])
    notes(s[14], "End on the tagline. Say it slowly and stop talking.")

    # --------------------------------------------------- 16 ACKNOWLEDGEMENTS
    strip_body(s[15])
    set_heading(s[15], "ACKNOWLEDGEMENTS")
    textbox(s[15], L, Y0, 7_900_000, 5_900_000, [
        {"t": "Team  [TEAM NAME]", "size": 22, "bold": True, "color": NAVY, "space": 12},
        {"t": "[MEMBER 1 — REG NO.]  —  Decentralized stack: ESP-NOW mesh, "
              "on-device auction, RF localization", "size": 17, "bullet": True, "space": 9},
        {"t": "[MEMBER 2 — REG NO.]  —  [ROLE]", "size": 17, "bullet": True, "space": 9},
        {"t": "[MEMBER 3 — REG NO.]  —  [ROLE]", "size": 17, "bullet": True, "space": 9},
        {"t": "[MEMBER 4 — REG NO.]  —  [ROLE]", "size": 17, "bullet": True, "space": 18},
        {"t": "Mentors", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "[MENTOR NAME(S) AND DEPARTMENT]", "size": 17, "bullet": True, "space": 18},
        {"t": "Tools", "size": 22, "bold": True, "color": NAVY, "space": 10},
        {"t": "PlatformIO · Arduino-ESP32 · Python 3 · MAVSDK-Python · YOLO",
         "size": 17, "space": 0},
    ])
    textbox(s[15], 9_300_000, Y0, 8_100_000, 5_900_000, [
        {"t": "References", "size": 22, "bold": True, "color": NAVY, "space": 12},
        {"t": "Gerkey, B. & Matarić, M. (2004). A Formal Analysis and Taxonomy of Task "
              "Allocation in Multi-Robot Systems. Int. Journal of Robotics Research.",
         "size": 16, "bullet": True, "space": 10},
        {"t": "Lagoudakis, M. et al. (2005). Auction-Based Multi-Robot Routing. "
              "Robotics: Science and Systems.", "size": 16, "bullet": True, "space": 10},
        {"t": "Espressif Systems. ESP-NOW Protocol Specification, ESP-IDF Programming Guide.",
         "size": 16, "bullet": True, "space": 10},
        {"t": "Stone, L. D. Theory of Optimal Search — Bayesian search theory, as applied "
              "in the recovery of USS Scorpion and Air France 447.",
         "size": 16, "bullet": True, "space": 10},
        {"t": "PX4 Autopilot. Hardware-in-the-Loop Simulation Documentation.",
         "size": 16, "bullet": True, "space": 0},
    ])
    notes(s[15], "Thank the mentors by name.")

    prs.save(DECK)
    print("saved", DECK, "-", len(s), "slides")


if __name__ == "__main__":
    main()
